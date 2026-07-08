"""Genera la presentazione PowerPoint di ClosetAI.

Output: ``docs/presentation.pptx`` (~16 slide per ~10-12 min, italiano).

Target: corso magistrale "Design per l'innovazione sostenibile", pre-esame.

**Narrazione**: la presentazione segue una *storia* — il ciclo di vita di un
capo nel guardaroba di Marco — e a ogni tappa mostra due lenti:

    👤  LATO UTENTE   → cosa vede e fa la persona
    ⚙️  LATO TECNICO  → cosa succede dietro, con la "natura" della tecnologia

Le quattro nature dell'AI/logica sono codificate a colori e ricorrono in
tutta la presentazione:

    🟦  AI PRE-ADDESTRATA   (Fashion-CLIP — riconosce i capi)
    🟩  AI ALLENATA DA NOI  (rete stato, gap analysis)
    🟪  AI GENERATIVA       (LLM/diffusion — descrizioni, coach, try-on)
    🟨  REGOLE / TABELLE    (wear log, cost-per-wear, colore, CO₂)

Ogni slide porta anche delle **note oratore** (script della parte orale da
esporre al professore) — visibili in PowerPoint nella vista Presentatore o
nel riquadro Note.

Uso::

    uv run python scripts/generate_presentation.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent.parent
SCREENSHOTS_DIR = ROOT / "docs" / "screenshots"
OUTPUT_PATH = ROOT / "docs" / "presentation.pptx"

# Palette base.
BG = RGBColor(0xF5, 0xF7, 0xFB)
INK = RGBColor(0x1A, 0x1D, 0x27)
MUTED = RGBColor(0x6B, 0x71, 0x82)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDC, 0xDF, 0xE8)
DANGER = RGBColor(0xC9, 0x4A, 0x5C)

# Le quattro nature (+ soft per gli sfondi).
PRE = RGBColor(0x4A, 0x6D, 0xDC)        # 🟦 AI pre-addestrata
PRE_SOFT = RGBColor(0xE0, 0xE8, 0xFA)
OWN = RGBColor(0x2F, 0x8F, 0x6E)        # 🟩 AI allenata da noi
OWN_SOFT = RGBColor(0xDC, 0xEF, 0xE7)
GEN = RGBColor(0x7A, 0x4F, 0xC4)        # 🟪 AI generativa
GEN_SOFT = RGBColor(0xEC, 0xE4, 0xF7)
RULE = RGBColor(0xC8, 0x8B, 0x2E)       # 🟨 regole / tabelle
RULE_SOFT = RGBColor(0xF6, 0xEC, 0xD8)

# Alias retro-compatibili (alcune slide li usano).
ACCENT = PRE
ACCENT_SOFT = PRE_SOFT
GREEN = OWN
GREEN_SOFT = OWN_SOFT
WARN = RULE

CATEGORY = {
    "pre": ("🟦 AI PRE-ADDESTRATA", PRE, PRE_SOFT),
    "own": ("🟩 AI ALLENATA DA NOI", OWN, OWN_SOFT),
    "gen": ("🟪 AI GENERATIVA", GEN, GEN_SOFT),
    "rule": ("🟨 REGOLE / TABELLE", RULE, RULE_SOFT),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_PAGES = 15


# ============================================================================
# Helpers di basso livello
# ============================================================================


def _set_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False,
                 color=INK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _add_card(slide, left, top, width, height, *, fill=PANEL, line=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    shape.text_frame.margin_left = Inches(0.18)
    shape.text_frame.margin_right = Inches(0.18)
    shape.text_frame.margin_top = Inches(0.14)
    shape.text_frame.margin_bottom = Inches(0.14)
    shape.text_frame.text = ""
    return shape


def _set_card_text(shape, lines: list[tuple[str, dict]]) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(opts.get("space_after", 4))
        run = p.add_run()
        run.text = text
        run.font.size = Pt(opts.get("size", 16))
        run.font.bold = opts.get("bold", False)
        run.font.italic = opts.get("italic", False)
        run.font.color.rgb = opts.get("color", INK)
        run.font.name = "Calibri"


def _add_slide_title(slide, title: str, kicker: str | None = None,
                     subtitle: str | None = None) -> None:
    y = Inches(0.42)
    if kicker:
        _add_textbox(slide, Inches(0.6), y, Inches(12), Inches(0.32),
                     kicker.upper(), size=11, bold=True, color=PRE)
        y += Inches(0.38)
    _add_textbox(slide, Inches(0.6), y, Inches(12), Inches(0.65),
                 title, size=30, bold=True, color=INK)
    if subtitle:
        _add_textbox(slide, Inches(0.6), y + Inches(0.66), Inches(12.2), Inches(0.5),
                     subtitle, size=14, color=MUTED, italic=True)
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.68), Inches(0.8), Emu(30000))
    underline.fill.solid()
    underline.fill.fore_color.rgb = PRE
    underline.line.fill.background()


def _add_footer(slide, page_num: int) -> None:
    _add_textbox(slide, Inches(0.6), Inches(7.12), Inches(11), Inches(0.3),
                 "ClosetAI · Design per l'innovazione sostenibile · pre-esame",
                 size=9, color=MUTED)
    _add_textbox(slide, Inches(11.8), Inches(7.12), Inches(1), Inches(0.3),
                 f"{page_num} / {TOTAL_PAGES}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def _add_screenshot_or_placeholder(slide, left, top, width, height,
                                   filename: str, caption: str):
    img_path = SCREENSHOTS_DIR / filename
    if img_path.is_file():
        slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
        return
    ph = _add_card(slide, left, top, width, height, fill=PRE_SOFT, line=PRE)
    _set_card_text(ph, [
        ("[ screenshot ]", {"size": 13, "bold": True, "color": PRE, "align": PP_ALIGN.CENTER, "space_after": 6}),
        (filename, {"size": 10, "color": MUTED, "align": PP_ALIGN.CENTER, "space_after": 4}),
        (caption, {"size": 10, "color": MUTED, "italic": True, "align": PP_ALIGN.CENTER}),
    ])


def _new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, BG)
    return slide


def _add_notes(slide, text: str) -> None:
    """Testo da esporre a voce, letto dal presentatore (non proiettato)."""
    slide.notes_slide.notes_text_frame.text = text


# ============================================================================
# Helper per le slide-tappa (doppia lente utente / tecnico)
# ============================================================================


def _tappa(prs, page, *, n, kicker, story_title, subtitle,
           user_lines, tech_categories, tech_lines,
           screenshot=None, screenshot_caption="", notes=""):
    """Slide a due colonne: 👤 lato utente | ⚙️ lato tecnico (con badge categorie)."""
    slide = _new_slide(prs)
    _add_slide_title(slide, story_title, kicker=f"Tappa {n} · {kicker}", subtitle=subtitle)

    col_y = Inches(2.0)
    col_h = Inches(4.8)
    left_x = Inches(0.6)
    right_x = Inches(6.95)
    col_w = Inches(5.95)

    # ── colonna UTENTE ──
    _add_textbox(slide, left_x, col_y, col_w, Inches(0.35),
                 "👤  LATO UTENTE — cosa vede e fa", size=12, bold=True, color=INK)
    if screenshot:
        _add_screenshot_or_placeholder(slide, left_x, col_y + Inches(0.45),
                                       col_w, Inches(4.2), screenshot, screenshot_caption)
    else:
        ucard = _add_card(slide, left_x, col_y + Inches(0.45), col_w, Inches(4.1),
                          fill=PANEL, line=BORDER)
        _set_card_text(ucard, [(t, o) for t, o in user_lines])

    # ── colonna TECNICO ──
    _add_textbox(slide, right_x, col_y, col_w, Inches(0.35),
                 "⚙️  LATO TECNICO — cosa succede dietro", size=12, bold=True, color=INK)
    tcard = _add_card(slide, right_x, col_y + Inches(0.45), col_w, col_h - Inches(0.65),
                      fill=PANEL, line=BORDER)
    lines: list[tuple[str, dict]] = []
    for cat_key in tech_categories:
        label, color, _ = CATEGORY[cat_key]
        lines.append((label, {"size": 12, "bold": True, "color": color, "space_after": 2}))
    lines.append(("", {"size": 4, "space_after": 2}))
    for t, o in tech_lines:
        lines.append((t, o))
    _set_card_text(tcard, lines)

    _add_footer(slide, page)
    if notes:
        _add_notes(slide, notes)
    return slide


# ============================================================================
# Slide
# ============================================================================


def slide_01_title(prs, page):
    slide = _new_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = PRE
    band.line.fill.background()

    _add_textbox(slide, Inches(0.9), Inches(1.5), Inches(11), Inches(0.5),
                 "PRE-ESAME · PRODOTTO AI ECO-SOSTENIBILE", size=12, bold=True, color=PRE)
    _add_textbox(slide, Inches(0.9), Inches(2.1), Inches(11), Inches(1.4),
                 "ClosetAI", size=80, bold=True, color=INK)
    _add_textbox(slide, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6),
                 "Il guardaroba intelligente che ti aiuta a vestire meglio comprando meno.",
                 size=22, color=OWN, italic=True)
    _add_textbox(slide, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.9),
                 "Una storia in sei tappe — dalla foto di un capo al suo fine vita — "
                 "raccontata da due punti di vista: chi la usa e cosa c'è dietro.",
                 size=14, color=MUTED)
    _add_textbox(slide, Inches(0.9), Inches(5.9), Inches(11), Inches(0.4),
                 "Design per l'innovazione sostenibile — Università di Pisa", size=15, color=MUTED)
    _add_textbox(slide, Inches(0.9), Inches(6.3), Inches(11), Inches(0.4),
                 "Marco Di Benedetto · A.A. 2025/2026", size=14, color=MUTED)
    _ = page
    _add_notes(slide,
        "Buongiorno, presento ClosetAI: un guardaroba intelligente che aiuta a "
        "vestire meglio comprando meno. Il prototipo copre l'intero ciclo di vita "
        "di un capo — dalla foto al fine vita — e lo racconto come una storia in "
        "sei tappe. A ogni tappa mostro due cose: cosa vede e fa l'utente, e cosa "
        "succede dietro le quinte a livello tecnico. È il filo conduttore di tutta "
        "la presentazione, quindi lo terrò esplicito slide per slide.")


def slide_02_problem(prs, page):
    slide = _new_slide(prs)
    _add_slide_title(slide, "Compriamo troppi vestiti.", kicker="01 · La sfida",
                     subtitle="L'industria della moda è una delle più impattanti del pianeta.")
    kpis = [
        ("10%", "delle emissioni\nglobali di CO₂", PRE),
        ("5 mln", "tonnellate di vestiti\nbuttati in Europa /anno", DANGER),
        ("7×", "utilizzi medi di un\ncapo fast fashion", RULE),
        ("40%", "del guardaroba medio\nmai indossato", DANGER),
    ]
    for i, (num, cap, color) in enumerate(kpis):
        x = Inches(0.6 + i * 3.1)
        card = _add_card(slide, x, Inches(2.3), Inches(2.9), Inches(2.0))
        _set_card_text(card, [
            (num, {"size": 44, "bold": True, "color": color, "align": PP_ALIGN.CENTER, "space_after": 6}),
            (cap, {"size": 12, "color": MUTED, "align": PP_ALIGN.CENTER}),
        ])
    _add_textbox(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
                 "Tre comportamenti, a casa nostra:", size=18, bold=True)
    _set_card_text(
        _add_card(slide, Inches(0.6), Inches(5.4), Inches(12.3), Inches(1.4), fill=PANEL, line=BORDER),
        [
            ("•  Compriamo d'impulso — non sappiamo davvero cosa abbiamo già", {"size": 14, "space_after": 4}),
            ("•  Capi 'fantasma' — vestiti dimenticati nell'armadio, mai indossati", {"size": 14, "space_after": 4}),
            ("•  Buttiamo invece di riparare, scambiare o donare", {"size": 14}),
        ])
    _add_footer(slide, page)
    _add_notes(slide,
        "Il punto di partenza è la scala del problema: la moda pesa per circa il "
        "10% delle emissioni globali di CO₂, e in Europa buttiamo 5 milioni di "
        "tonnellate di vestiti l'anno, spesso quasi nuovi — un capo fast fashion "
        "viene indossato in media solo 7 volte. Ma il problema non è solo "
        "industriale, è anche domestico: compriamo d'impulso senza sapere cosa "
        "abbiamo già, dimentichiamo capi in armadio, e buttiamo invece di "
        "riparare o donare. ClosetAI attacca proprio questi tre comportamenti.")


def slide_03_story_intro(prs, page):
    slide = _new_slide(prs)
    _add_slide_title(slide, "Seguiamo un capo, dalla foto al fine vita.",
                     kicker="02 · La storia + la legenda",
                     subtitle="A ogni tappa: cosa fa l'utente e quale tecnologia lavora dietro.")

    # Le 6 tappe in fila
    tappe = ["📷\nLo fotografo", "✓\nLo indosso", "👗\nCosa metto?",
             "🛠️\nSi è rovinato", "🧩\nMi serve altro?", "♻️\nLo lascio andare"]
    tw = Inches(1.95)
    for i, t in enumerate(tappe):
        x = Inches(0.6) + (tw + Inches(0.12)) * i
        card = _add_card(slide, x, Inches(2.05), tw, Inches(1.5), fill=PANEL, line=BORDER)
        _set_card_text(card, [
            (f"{i+1}", {"size": 13, "bold": True, "color": PRE, "align": PP_ALIGN.CENTER, "space_after": 2}),
            (t, {"size": 12, "color": INK, "align": PP_ALIGN.CENTER}),
        ])
        if i < len(tappe) - 1:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + tw + Inches(0.005),
                                        Inches(2.65), Inches(0.11), Inches(0.3))
            ar.fill.solid()
            ar.fill.fore_color.rgb = PRE
            ar.line.fill.background()

    # Legenda 4 nature
    _add_textbox(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
                 "Le quattro nature della tecnologia (codice colore in tutta la presentazione):",
                 size=14, bold=True, color=INK)
    legend = [
        ("pre", "Riconosce, da un modello già addestrato da altri.\nEs. Fashion-CLIP riconosce i capi."),
        ("own", "Una rete che abbiamo addestrato noi.\nEs. stato del capo, vuoti del guardaroba."),
        ("gen", "Produce contenuti nuovi (testi, immagini).\nEs. descrizioni, coach, try-on virtuale."),
        ("rule", "Logica esplicita e trasparente, niente ML.\nEs. cost-per-wear, tabella CO₂."),
    ]
    lw = Inches(3.0)
    for i, (key, desc) in enumerate(legend):
        label, color, soft = CATEGORY[key]
        x = Inches(0.6) + (lw + Inches(0.1)) * i
        card = _add_card(slide, x, Inches(4.5), lw, Inches(2.2), fill=soft, line=color)
        _set_card_text(card, [
            (label, {"size": 12, "bold": True, "color": color, "align": PP_ALIGN.CENTER, "space_after": 8}),
            (desc, {"size": 11, "color": MUTED, "align": PP_ALIGN.CENTER}),
        ])
    _add_footer(slide, page)
    _add_notes(slide,
        "Il capo che seguiamo attraversa sei momenti: lo fotografo, lo indosso, "
        "chiedo cosa mettermi, si rovina, mi chiedo se mi serve altro, e infine "
        "lo lascio andare. Questa è la mappa di tutta la presentazione. "
        "Prima di entrare nel dettaglio, fisso la legenda: useremo sempre lo "
        "stesso codice colore per capire CHE TIPO di intelligenza sta lavorando — "
        "blu per un modello pre-addestrato da altri, verde per le reti che "
        "abbiamo addestrato noi, viola per l'AI generativa che crea contenuti "
        "nuovi, giallo per la logica a regole, esplicita e trasparente. Tenetela "
        "a mente: la ritroverete a ogni tappa.")


def slide_04_tappa_riconosci(prs, page):
    _tappa(
        prs, page, n=1, kicker="«L'ho appena comprato»",
        story_title="Lo fotografo, l'app lo riconosce.",
        subtitle="Niente moduli da compilare: basta una foto.",
        screenshot="01-home.png", screenshot_caption="il guardaroba digitale",
        user_lines=[
            ("Marco fotografa un capo nuovo.", {"size": 15, "bold": True, "space_after": 10}),
            ("In un attimo lo ritrova nel guardaroba", {"size": 13, "color": MUTED, "space_after": 4}),
            ("digitale, già etichettato con", {"size": 13, "color": MUTED, "space_after": 4}),
            ("categoria e colore — senza scrivere", {"size": 13, "color": MUTED, "space_after": 4}),
            ("nulla a mano.", {"size": 13, "color": MUTED}),
        ],
        tech_categories=["pre"],
        tech_lines=[
            ("Fashion-CLIP", {"size": 15, "bold": True, "space_after": 4}),
            ("Rete neurale pre-addestrata su 700.000 immagini", {"size": 12, "color": MUTED, "space_after": 2}),
            ("di moda. Riconosce la categoria (t-shirt, giacca…)", {"size": 12, "color": MUTED, "space_after": 2}),
            ("e produce un 'embedding' del capo.", {"size": 12, "color": MUTED, "space_after": 10}),
            ("Perché pre-addestrata e non nostra?", {"size": 12, "bold": True, "color": PRE, "space_after": 2}),
            ("Riaddestrarla servirebbero decine di migliaia di", {"size": 11, "color": MUTED, "space_after": 2}),
            ("foto etichettate: la usiamo già pronta.", {"size": 11, "color": MUTED}),
        ],
        notes=(
            "Prima tappa: Marco compra un capo e lo fotografa. Niente form da "
            "compilare — la foto passa a Fashion-CLIP, un modello pre-addestrato "
            "su circa 700 mila immagini di moda, che riconosce categoria e colore "
            "e produce anche un embedding visivo, un vettore che userò più avanti "
            "per la varietà degli outfit. Perché non l'abbiamo addestrato noi? "
            "Perché sarebbe irrazionale: servirebbero decine di migliaia di foto "
            "etichettate per un problema che qualcuno ha già risolto meglio di "
            "quanto potremmo fare noi con le risorse di un progetto universitario."
        ),
    )


def slide_05_tappa_indosso(prs, page):
    _tappa(
        prs, page, n=2, kicker="«Oggi l'ho messo»",
        story_title="Lo indosso, l'app tiene il conto.",
        subtitle="Un solo tap trasforma l'armadio in dati utili.",
        user_lines=[
            ("Marco tocca '✓ indossato oggi'.", {"size": 15, "bold": True, "space_after": 10}),
            ("Col tempo scopre quanto gli costa", {"size": 13, "color": MUTED, "space_after": 4}),
            ("davvero ogni capo (cost-per-wear) e", {"size": 13, "color": MUTED, "space_after": 4}),
            ("quali sono i suoi capi 'fantasma' —", {"size": 13, "color": MUTED, "space_after": 4}),
            ("mai indossati da mesi.", {"size": 13, "color": MUTED}),
        ],
        tech_categories=["rule"],
        tech_lines=[
            ("Wear log + due calcoli", {"size": 15, "bold": True, "space_after": 6}),
            ("• cost-per-wear = prezzo ÷ numero di utilizzi", {"size": 12, "color": MUTED, "space_after": 4}),
            ("• capo fantasma = 0 utilizzi dopo N giorni", {"size": 12, "color": MUTED, "space_after": 10}),
            ("Perché regole e non una rete?", {"size": 12, "bold": True, "color": RULE, "space_after": 2}),
            ("Sono formule esatte e trasparenti: una divisione", {"size": 11, "color": MUTED, "space_after": 2}),
            ("e una soglia. Una rete qui sarebbe inutile e meno", {"size": 11, "color": MUTED, "space_after": 2}),
            ("affidabile di un conto che torna sempre.", {"size": 11, "color": MUTED}),
        ],
        notes=(
            "Seconda tappa: l'uso reale. Un tap su 'indossato oggi' registra "
            "l'evento, e da lì calcoliamo due metriche semplici ma potenti: il "
            "cost-per-wear, cioè prezzo diviso numero di utilizzi — quanto ti è "
            "davvero costato quel capo — e i capi 'fantasma', quelli mai indossati "
            "dopo N giorni. Qui uso deliberatamente delle regole e non una rete "
            "neurale: sono formule esatte, trasparenti, verificabili a mano. "
            "Un modello ML introdurrebbe incertezza dove non ne serve: la "
            "matematica qui è già la risposta giusta."
        ),
    )


def slide_06_tappa_outfit(prs, page):
    _tappa(
        prs, page, n=3, kicker="«Cosa mi metto oggi?»",
        story_title="L'app propone l'outfit del giorno.",
        subtitle="Dai capi che già possiede, non da un catalogo da comprare.",
        screenshot="04-today.png", screenshot_caption="proposte outfit + meteo",
        user_lines=[
            ("Marco chiede aiuto la mattina.", {"size": 15, "bold": True, "space_after": 10}),
            ("Riceve 3 proposte composte SOLO", {"size": 13, "color": MUTED, "space_after": 4}),
            ("con i suoi capi, adatte al meteo del", {"size": 13, "color": MUTED, "space_after": 4}),
            ("giorno, e dà un like a quelle che", {"size": 13, "color": MUTED, "space_after": 4}),
            ("preferisce.", {"size": 13, "color": MUTED}),
        ],
        tech_categories=["rule", "pre"],
        tech_lines=[
            ("Recommender: regole + embedding + meteo", {"size": 14, "bold": True, "space_after": 6}),
            ("• 🟨 regole di colore (complementari, neutri…)", {"size": 12, "color": MUTED, "space_after": 3}),
            ("• 🟨 filtro meteo via API Open-Meteo", {"size": 12, "color": MUTED, "space_after": 3}),
            ("• 🟦 l'embedding di Fashion-CLIP dà varietà", {"size": 12, "color": MUTED, "space_after": 10}),
            ("Un mix: la logica estetica è scritta a mano,", {"size": 11, "italic": True, "color": MUTED, "space_after": 2}),
            ("ma sfrutta la 'percezione' del modello pre-addestrato.", {"size": 11, "italic": True, "color": MUTED}),
        ],
        notes=(
            "Terza tappa: 'cosa mi metto oggi?'. Il recommender combina tre "
            "ingredienti: regole di compatibilità cromatica scritte a mano — "
            "colori complementari, presenza di neutri — un filtro sul meteo del "
            "giorno via l'API Open-Meteo, e un bonus di varietà che sfrutta "
            "l'embedding di Fashion-CLIP per non riproporre sempre le stesse "
            "combinazioni. È un esempio di mix consapevole: la logica estetica "
            "resta esplicita e ispezionabile, ma appoggiata sulla percezione del "
            "modello pre-addestrato per la parte che sarebbe troppo costoso "
            "codificare a mano."
        ),
    )


def slide_07_tappa_rovinato(prs, page):
    _tappa(
        prs, page, n=4, kicker="«Si è rovinato»",
        story_title="L'app riconosce lo stato dalla foto.",
        subtitle="La nostra prima rete neurale, e la sua validazione onesta.",
        user_lines=[
            ("Marco fotografa un capo logoro.", {"size": 15, "bold": True, "space_after": 10}),
            ("L'app gli dice se è in buono stato,", {"size": 13, "color": MUTED, "space_after": 4}),
            ("usurato o danneggiato — e da lì", {"size": 13, "color": MUTED, "space_after": 4}),
            ("parte il suggerimento dell'azione", {"size": 13, "color": MUTED, "space_after": 4}),
            ("circolare migliore (tappa 6).", {"size": 13, "color": MUTED}),
        ],
        tech_categories=["own"],
        tech_lines=[
            ("Rete dello STATO (addestrata da noi)", {"size": 15, "bold": True, "space_after": 6}),
            ("Fashion-CLIP (frozen) → MLP 512→256→128→3", {"size": 12, "color": MUTED, "space_after": 8}),
            ("Validazione onesta:", {"size": 12, "bold": True, "color": OWN, "space_after": 2}),
            ("dati sintetici → 96% (sovrastima) · su foto reali", {"size": 11, "color": MUTED, "space_after": 2}),
            ("con difetti veri (dataset COCO) crolla al 60%.", {"size": 11, "color": MUTED, "space_after": 4}),
            ("Causa: confine artificiale fra due classi troppo", {"size": 11, "color": MUTED, "space_after": 2}),
            ("simili. Fuse in 3 classi oneste → 94% reale.", {"size": 11, "color": MUTED, "space_after": 8}),
            ("Un errore trovato e corretto, non nascosto.", {"size": 11, "italic": True, "color": MUTED}),
        ],
        notes=(
            "Quarta tappa, e la parte di cui vado più fiero dal punto di vista "
            "metodologico. Una foto entra in Fashion-CLIP, che dà l'embedding, e "
            "sopra quell'embedding abbiamo addestrato noi un piccolo MLP che "
            "predice lo stato di conservazione. Sul dataset sintetico arrivavamo "
            "al 96%, un numero bellissimo ma sospetto. Appena l'ho testato su "
            "foto reali con difetti veri, l'accuracy è crollata al 60% — un "
            "classico domain gap. Invece di nasconderlo, l'ho investigato: la "
            "confusione era quasi tutta fra due classi, 'nuovo' e 'buono', "
            "indistinguibili su foto reali. Le ho fuse in un'unica classe e "
            "l'accuracy onesta è risalita al 94%, con il danneggiato riconosciuto "
            "quasi perfettamente. Un promemoria: le vecchie feature legate al "
            "tutorial di riparazione generato via VLM sono state rimosse dal "
            "prodotto in una revisione successiva, perché aggiungevano "
            "complessità senza un beneficio proporzionato — coerenti col "
            "principio MVP first che ripeto più avanti."
        ),
    )


def slide_08_tappa_serve(prs, page):
    _tappa(
        prs, page, n=5, kicker="«Mi serve davvero qualcosa?»",
        story_title="L'app legge l'armadio e trova i vuoti.",
        subtitle="Una rete che ragiona sull'insieme, non sul singolo capo.",
        user_lines=[
            ("Prima di comprare, Marco controlla.", {"size": 15, "bold": True, "space_after": 10}),
            ("L'app gli dice se l'armadio è", {"size": 13, "color": MUTED, "space_after": 4}),
            ("equilibrato o se ha dei vuoti reali", {"size": 13, "color": MUTED, "space_after": 4}),
            ("(\"manca una giacca\", \"troppe", {"size": 13, "color": MUTED, "space_after": 4}),
            ("t-shirt\") e consiglia, se serve,", {"size": 13, "color": MUTED, "space_after": 4}),
            ("di cercarlo second-hand.", {"size": 13, "color": MUTED}),
        ],
        tech_categories=["own"],
        tech_lines=[
            ("Rete di gap analysis (addestrata da noi)", {"size": 14, "bold": True, "space_after": 6}),
            ("Non guarda le foto: guarda l'inventario prodotto", {"size": 12, "color": MUTED, "space_after": 2}),
            ("da Fashion-CLIP. Input: quanti capi per categoria,", {"size": 12, "color": MUTED, "space_after": 2}),
            ("colori, stagioni, frequenza d'uso →", {"size": 12, "color": MUTED, "space_after": 2}),
            ("Output: i vuoti funzionali del guardaroba.", {"size": 12, "color": MUTED, "space_after": 10}),
            ("Perché una rete e non regole?", {"size": 12, "bold": True, "color": OWN, "space_after": 2}),
            ("Cattura combinazioni sfumate fra molti fattori,", {"size": 11, "color": MUTED, "space_after": 2}),
            ("e migliorerà col feedback reale degli utenti.", {"size": 11, "color": MUTED}),
        ],
        notes=(
            "Quinta tappa: prima di comprare, conviene guardare l'armadio nel "
            "suo insieme. Qui cambia livello: non guardiamo più una foto ma i "
            "dati aggregati che Fashion-CLIP ha già prodotto capo per capo — "
            "quanti per categoria, quali colori, quanto vengono usati. Una "
            "seconda rete che abbiamo addestrato noi, multi-label, impara a "
            "riconoscere i vuoti funzionali: manca una giacca, ci sono troppe "
            "t-shirt. È addestrata su un dataset sintetico con etichette da "
            "regole esperte più rumore, così la rete impara le soglie invece di "
            "copiarle: un limite dichiarato, che miglioreremmo con feedback "
            "reale degli utenti."
        ),
    )


def slide_09_tappa_circolare(prs, page):
    _tappa(
        prs, page, n=6, kicker="«Non lo uso più»",
        story_title="L'app lo accompagna a una seconda vita.",
        subtitle="E misura quanta CO₂ hai risparmiato così.",
        user_lines=[
            ("Marco decide il destino del capo.", {"size": 15, "bold": True, "space_after": 10}),
            ("L'app propone l'azione migliore —", {"size": 13, "color": MUTED, "space_after": 4}),
            ("riparare, scambiare, vendere,", {"size": 13, "color": MUTED, "space_after": 4}),
            ("donare o riciclare — e gli mostra", {"size": 13, "color": MUTED, "space_after": 4}),
            ("quanti kg di CO₂ ha evitato.", {"size": 13, "color": MUTED}),
        ],
        tech_categories=["rule"],
        tech_lines=[
            ("Tabella CO₂ (Ellen MacArthur Foundation)", {"size": 14, "bold": True, "space_after": 6}),
            ("impatto del capo × % evitata dall'azione:", {"size": 12, "color": MUTED, "space_after": 4}),
            ("• donare / vendere / scambiare → 100%", {"size": 12, "color": OWN, "space_after": 3}),
            ("• riparare → 70%   • riciclare → 30%", {"size": 12, "color": OWN, "space_after": 10}),
            ("Perché una tabella e non una rete?", {"size": 12, "bold": True, "color": RULE, "space_after": 2}),
            ("L'impatto CO₂ è già stimato da database LCA", {"size": 11, "color": MUTED, "space_after": 2}),
            ("autorevoli: meglio un numero citabile e trasparente", {"size": 11, "color": MUTED, "space_after": 2}),
            ("che una scatola nera.", {"size": 11, "color": MUTED}),
        ],
        notes=(
            "Sesta e ultima tappa del ciclo di vita: il capo non serve più a "
            "Marco. In base allo stato diagnosticato alla tappa 4, l'app "
            "suggerisce l'azione più coerente con la gerarchia dell'economia "
            "circolare — riparare prima di riciclare, rivendere o donare prima "
            "di buttare — e stima i kg di CO₂ evitati con una tabella per "
            "categoria di capo, presa da fonti LCA autorevoli come la Ellen "
            "MacArthur Foundation. Anche qui niente rete neurale: la relazione "
            "materiale-impatto è già nota, e una rete darebbe solo falsa "
            "precisione dove serve invece un numero trasparente e citabile."
        ),
    )


def slide_10_dashboard(prs, page):
    slide = _new_slide(prs)
    _add_slide_title(slide, "Il premio: vedere il proprio impatto.",
                     kicker="03 · Il payoff per l'utente",
                     subtitle="Tutte le tappe convergono in numeri concreti e motivanti.")
    _add_screenshot_or_placeholder(slide, Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.7),
                                   "05-dashboard.png", "dashboard impatto")
    card = _add_card(slide, Inches(8.9), Inches(2.0), Inches(4.0), Inches(4.7),
                     fill=OWN_SOFT, line=OWN)
    _set_card_text(card, [
        ("Cosa vede Marco:", {"size": 14, "bold": True, "color": OWN, "space_after": 8}),
        ("→ CO₂ evitata, in km auto", {"size": 13, "color": INK, "space_after": 6}),
        ("→ capi salvati dalla discarica", {"size": 13, "color": INK, "space_after": 6}),
        ("→ cost-per-wear del guardaroba", {"size": 13, "color": INK, "space_after": 6}),
        ("→ vuoti da colmare (o no!)", {"size": 13, "color": INK, "space_after": 6}),
        ("→ coach AI con un consiglio", {"size": 13, "color": INK, "space_after": 12}),
        ("La sostenibilità diventa", {"size": 12, "italic": True, "color": MUTED, "space_after": 2}),
        ("misurabile e personale.", {"size": 12, "italic": True, "color": MUTED}),
    ])
    _add_footer(slide, page)
    _add_notes(slide,
        "Tutte e sei le tappe convergono qui, nella dashboard impatto: CO₂ "
        "evitata tradotta in km d'auto risparmiati per renderla tangibile, capi "
        "salvati dalla discarica, il cost-per-wear di tutto il guardaroba, i "
        "vuoti funzionali da colmare o no, e un coach — questo sì AI generativa, "
        "un LLM via litellm — che scrive un consiglio personalizzato basato sui "
        "dati reali dell'utente. È il momento in cui la sostenibilità smette di "
        "essere un'affermazione astratta e diventa un numero personale e "
        "motivante.")


def slide_11_pipeline(prs, page):
    slide = _new_slide(prs)
    _add_slide_title(slide, "Il quadro tecnico, in un colpo d'occhio.",
                     kicker="04 · Sotto il cofano",
                     subtitle="Le stesse tappe, ora come pipeline su tre livelli. Colori = le 4 nature.")

    box_w = Inches(4.0)
    gap = Inches(0.17)

    def row(y, boxes):
        for i, (title, kind, out, color, soft) in enumerate(boxes):
            x = Inches(0.6) + (box_w + gap) * i
            card = _add_card(slide, x, y, box_w, Inches(1.2), fill=soft, line=color)
            _set_card_text(card, [
                (title, {"size": 12.5, "bold": True, "color": INK, "space_after": 1}),
                (kind, {"size": 9, "italic": True, "color": color, "space_after": 4}),
                (out, {"size": 10, "color": MUTED}),
            ])

    _add_textbox(slide, Inches(0.6), Inches(1.78), Inches(12), Inches(0.3),
                 "📷  PER OGNI CAPO · dalla foto", size=11, bold=True, color=INK)
    row(Inches(2.08), [
        ("Fashion-CLIP", "🟦 pre-addestrata", "→ categoria + colore (+ embedding)", PRE, PRE_SOFT),
        ("Rete stato del capo", "🟩 nostra", "→ buono / usurato / danneggiato", OWN, OWN_SOFT),
        ("Tabella CO₂", "🟨 regole", "→ impatto di produzione", RULE, RULE_SOFT),
    ])
    _add_textbox(slide, Inches(0.6), Inches(3.48), Inches(12), Inches(0.3),
                 "✓  USO NEL TEMPO · un tap 'indossato oggi'", size=11, bold=True, color=INK)
    row(Inches(3.78), [
        ("Wear log", "🟨 regole", "→ quante volte l'hai indossato", RULE, RULE_SOFT),
        ("Cost-per-wear", "🟨 calcolo", "→ prezzo ÷ n. utilizzi", RULE, RULE_SOFT),
        ("Capi fantasma", "🟨 soglia", "→ mai indossati > N giorni", RULE, RULE_SOFT),
    ])
    _add_textbox(slide, Inches(0.6), Inches(5.18), Inches(12), Inches(0.3),
                 "🧩  TUTTO IL GUARDAROBA · sull'insieme", size=11, bold=True, color=INK)
    row(Inches(5.48), [
        ("Outfit recommender", "🟨 regole + 🟦 + 🌦", "→ \"cosa metto oggi?\"", PRE, PRE_SOFT),
        ("Rete gap analysis", "🟩 nostra", "→ vuoti: \"manca una giacca\"", OWN, OWN_SOFT),
        ("Somma impatti evitati", "🟨 calcolo", "→ CO₂ totale → dashboard", RULE, RULE_SOFT),
    ])
    _add_footer(slide, page)
    _add_notes(slide,
        "Riassumo tutto in un'unica vista, su tre livelli. Per ogni singola "
        "foto lavorano tre motori: Fashion-CLIP pre-addestrato, la nostra rete "
        "dello stato, la tabella CO₂. Nel tempo, ogni 'indossato oggi' alimenta "
        "wear log, cost-per-wear e capi fantasma, tutta logica a regole. E "
        "sull'intero guardaroba lavorano l'outfit recommender e la rete di gap "
        "analysis, fino alla somma degli impatti evitati che arriva in "
        "dashboard. Il messaggio di questa slide è che nessun livello dipende "
        "da AI generativa nel suo cuore: quella la usiamo altrove, come "
        "supporto opzionale, non come collante indispensabile della pipeline.")


def slide_12_four_natures(prs, page):
    slide = _new_slide(prs)
    _add_slide_title(slide, "Quattro nature, ognuna dove ha senso.",
                     kicker="05 · La tesi tecnica",
                     subtitle="Non 'AI ovunque': lo strumento giusto per ogni problema.")
    cards = [
        ("pre", "AI PRE-ADDESTRATA", "Fashion-CLIP",
         "Quando il problema (riconoscere un capo) è\ngenerale e qualcun altro l'ha già risolto su\nmilioni di immagini. Non lo rifacciamo noi."),
        ("own", "AI ALLENATA DA NOI", "rete stato · gap analysis",
         "Quando il problema è specifico del nostro\ndominio e abbiamo (o generiamo) i dati per\ninsegnarlo. Qui dimostriamo il know-how ML."),
        ("gen", "AI GENERATIVA", "LLM · diffusion",
         "Quando serve creare contenuto nuovo e\npersonalizzato: descrizioni narrative,\ncoach di sostenibilità, prova virtuale."),
        ("rule", "REGOLE / TABELLE", "cost-per-wear · CO₂ · colore",
         "Quando la relazione è già nota, esatta e va\ntenuta trasparente. Una rete qui darebbe solo\nfalsa precisione e meno fiducia."),
    ]
    for i, (key, title, sub, desc) in enumerate(cards):
        col = i % 2
        rowi = i // 2
        x = Inches(0.6 + col * 6.35)
        y = Inches(2.05 + rowi * 2.45)
        _, color, soft = CATEGORY[key]
        card = _add_card(slide, x, y, Inches(6.1), Inches(2.25), fill=soft, line=color)
        _set_card_text(card, [
            (title, {"size": 15, "bold": True, "color": color, "space_after": 1}),
            (sub, {"size": 11, "italic": True, "color": MUTED, "space_after": 8}),
            (desc, {"size": 12, "color": INK}),
        ])
    _add_footer(slide, page)
    _add_notes(slide,
        "Questa è la tesi tecnica del progetto: non 'AI ovunque', ma lo "
        "strumento giusto per ogni problema, e la disciplina di sapere quando "
        "NON usare AI. Pre-addestrata quando il problema è generale e già "
        "risolto da altri su scala che non potremmo replicare. Allenata da noi "
        "quando il problema è specifico del nostro dominio — qui dimostriamo il "
        "know-how di machine learning vero e proprio. Generativa quando serve "
        "contenuto nuovo e personalizzato. Regole quando la relazione è già "
        "nota ed esatta. Una nota di onestà: avevamo prototipato anche un "
        "quarto approccio, un modello generativo che scriveva tutorial di "
        "riparazione dalla foto — funzionava, ma l'abbiamo rimosso in revisione "
        "perché il valore aggiunto non giustificava la complessità di due "
        "modelli vision-generativi extra da mantenere. Anche sapere cosa "
        "togliere è parte del metodo MVP first.")


def slide_13_feasibility(prs, page):
    slide = _new_slide(prs)
    _add_slide_title(slide, "Studio di fattibilità.",
                     kicker="06 · Costi, materiali, scalabilità",
                     subtitle="Si può davvero portarlo sul mercato?")
    cols = [
        ("💰 Costi", PRE, [
            "Hardware: laptop + telefono.",
            "Modelli AI: gratuiti / open-source.",
            "Cloud opzionale: ~€ 0,01 a utente/mese",
            "(grazie alla cache delle risposte).",
            "Specchio smart: ~€ 200 una tantum.",
        ]),
        ("🧱 Materiali", OWN, [
            "Tutto open-source, nessuna licenza:",
            "FastAPI, React, Fashion-CLIP, Qwen-VL,",
            "Stable Diffusion, scikit-learn, PyTorch.",
            "Piena trasparenza e nessun lock-in",
            "su un fornitore.",
        ]),
        ("📈 Scalabilità", RULE, [
            "MVP: single-user, locale.",
            "→ Cloud multi-utente (Postgres + S3).",
            "→ AI on-device per la privacy.",
            "→ Versione con marketplace",
            "second-hand integrato.",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        x = Inches(0.6 + i * 4.2)
        card = _add_card(slide, x, Inches(2.1), Inches(4.0), Inches(4.7), fill=PANEL, line=color)
        lines = [(title, {"size": 18, "bold": True, "color": color, "space_after": 10})]
        for it in items:
            lines.append((it, {"size": 12, "color": INK, "space_after": 6}))
        _set_card_text(card, lines)
    _add_footer(slide, page)
    _add_notes(slide,
        "Un progetto universitario deve anche reggere la domanda 'si può "
        "portarlo sul mercato?'. Sui costi: hardware minimo — laptop e "
        "telefono — modelli AI gratuiti o open-source, e un costo cloud "
        "trascurabile grazie alla cache delle risposte LLM; lo specchio smart "
        "opzionale costa sui 200 euro una tantum in componenti Raspberry Pi. "
        "Sui materiali: stack interamente open-source, FastAPI, React, "
        "Fashion-CLIP, Stable Diffusion, PyTorch — nessun lock-in su un "
        "fornitore. Sulla scalabilità: partiamo single-user locale per "
        "l'MVP, ma il percorso verso multi-utente su cloud, AI on-device per "
        "la privacy, e un marketplace second-hand integrato è già chiaro.")


def slide_14_limits(prs, page):
    slide = _new_slide(prs)
    _add_slide_title(slide, "Cosa il prototipo non fa (ancora).",
                     kicker="07 · Onestà intellettuale",
                     subtitle="E dove potrebbe andare in futuro.")
    cols = [
        ("Limitazioni di oggi", DANGER, [
            "I dataset delle reti che alleniamo noi sono sintetici: vanno validati su dati reali.",
            "Le stime CO₂ sono medie per categoria, non per materiale.",
            "Il try-on è un'illusione visiva, non un camerino digitale.",
            "Per ora è personale: niente account né condivisione.",
        ]),
        ("Estensioni naturali", OWN, [
            "Raccolta di foto reali per riaddestrare la rete dello stato.",
            "Marketplace second-hand integrato (Vinted, Wallapop).",
            "CO₂ per materiale (cotone vs lana vs poliestere).",
            "Specchio smart on-device, senza inviare foto in cloud.",
        ]),
    ]
    for i, (title, color, items) in enumerate(cols):
        x = Inches(0.6 + i * 6.3)
        card = _add_card(slide, x, Inches(2.1), Inches(6.1), Inches(4.7), fill=PANEL, line=color)
        lines = [(title, {"size": 18, "bold": True, "color": color, "space_after": 12})]
        for it in items:
            lines.append((f"•  {it}", {"size": 13, "color": INK, "space_after": 10}))
        _set_card_text(card, lines)
    _add_footer(slide, page)
    _add_notes(slide,
        "Chiudo con onestà intellettuale, che per me è parte integrante del "
        "metodo, non un vezzo finale. I dataset delle reti che alleniamo noi "
        "restano parzialmente sintetici e vanno validati su dati reali — "
        "l'ho mostrato concretamente con la rete dello stato. Le stime CO₂ "
        "sono medie per categoria, non ancora per materiale. Il try-on è "
        "un'illusione visiva convincente ma non un vero camerino digitale. "
        "Ed è personale: oggi niente account condivisi. Le estensioni naturali "
        "sono altrettanto chiare: più foto reali, marketplace integrato, CO₂ "
        "per materiale, e uno specchio smart che elabora tutto on-device senza "
        "mai inviare foto in cloud — coerente col principio di privacy by "
        "design che abbiamo seguito fin dall'inizio.")


def slide_15_closing(prs, page):
    slide = _new_slide(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), SLIDE_H)
    band.fill.solid()
    band.fill.fore_color.rgb = OWN
    band.line.fill.background()
    _add_textbox(slide, Inches(0.9), Inches(1.9), Inches(11), Inches(1.1),
                 "Grazie.", size=72, bold=True, color=INK)
    _add_textbox(slide, Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.7),
                 "Vesti meglio. Compra meno. Allunga la vita dei tuoi capi.",
                 size=22, color=OWN, italic=True)
    _add_textbox(slide, Inches(0.9), Inches(4.3), Inches(11), Inches(0.5),
                 "Domande?", size=24, color=PRE)
    box = _add_card(slide, Inches(0.9), Inches(5.2), Inches(11.6), Inches(1.5), fill=PANEL, line=BORDER)
    _set_card_text(box, [
        ("Il filo conduttore:", {"size": 13, "bold": True, "color": PRE, "space_after": 4}),
        ("una storia in 6 tappe, dove ogni problema ha lo strumento giusto —",
         {"size": 13, "color": INK, "space_after": 4}),
        ("🟦 modello pre-addestrato · 🟩 rete nostra · 🟪 AI generativa · 🟨 regole trasparenti.",
         {"size": 13, "color": MUTED}),
    ])
    _ = page
    _add_notes(slide,
        "Per chiudere: ClosetAI nasce da una storia in sei tappe, e in ogni "
        "tappa ho cercato lo strumento giusto per il problema — non l'AI più "
        "sofisticata possibile, ma quella più onesta. Un modello pre-addestrato "
        "dove il problema è già risolto, una rete nostra dove serve know-how "
        "specifico e l'abbiamo dimostrato con dati e validazione reali, l'AI "
        "generativa dove serve creare contenuto, le regole trasparenti dove la "
        "risposta è già nota. Grazie per l'attenzione, sono a disposizione per "
        "domande.")


# ============================================================================
# Build
# ============================================================================


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        slide_01_title,
        slide_02_problem,
        slide_03_story_intro,
        slide_04_tappa_riconosci,
        slide_05_tappa_indosso,
        slide_06_tappa_outfit,
        slide_07_tappa_rovinato,
        slide_08_tappa_serve,
        slide_09_tappa_circolare,
        slide_10_dashboard,
        slide_11_pipeline,
        slide_12_four_natures,
        slide_13_feasibility,
        slide_14_limits,
        slide_15_closing,
    ]
    for i, builder in enumerate(slides, start=1):
        builder(prs, i)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    out = build()
    expected = ("01-home.png", "04-today.png", "05-dashboard.png")
    found = sum(1 for n in expected if (SCREENSHOTS_DIR / n).is_file())
    print(f"==> Presentazione generata: {out}")
    print("    Slide totali: 15  (target: ~10-12 min)")
    print(f"    Screenshot usati: {found}/{len(expected)}")
    if found < len(expected):
        missing = [n for n in expected if not (SCREENSHOTS_DIR / n).is_file()]
        print(f"    Mancanti (placeholder): {', '.join(missing)}")
    sys.exit(0)
