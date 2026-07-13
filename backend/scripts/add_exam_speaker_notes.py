"""Add the official oral commentary to the generated ClosetAI deck.

The visual deck is built with artifact-tool.  This small post-processing step
uses python-pptx only to populate PowerPoint's native speaker-notes parts.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


NOTES = [
    """Apro con una distinzione importante: ClosetAI non è “AI applicata alla moda” in senso generico. È un prodotto che prova a cambiare una decisione quotidiana: usare più a lungo ciò che possediamo prima di comprare altro. Durante la presentazione mostrerò quattro segnali visibili nell’interfaccia—impatto verde, stato dei capi, gap del guardaroba e capi fantasma—e dirò sempre se dietro c’è un modello nostro, un modello pre-addestrato, un servizio generativo oppure una regola. Questa precisione è parte del progetto, non una nota a margine.""",
    """Per motivare il problema uso solo numeri che posso citare. UNEP stima che il tessile produca fra il 2 e l’8 per cento delle emissioni globali: presento un intervallo, non un falso numero esatto. L’EEA riporta che nel 2022 in UE si consumavano 19 chilogrammi di tessili per persona e se ne generavano 16 come rifiuto. Inoltre l’85 per cento dei rifiuti tessili domestici non veniva raccolto separatamente. Il punto non è colpevolizzare l’utente: è rendere più facile riuso, cura e seconda vita.""",
    """Questa è la tesi di prodotto. ClosetAI chiude un ciclo: fotografo il capo, registro quando lo indosso, riscopro ciò che non uso, valuto lo stato, decido se esiste davvero un gap e, quando serve, scelgo un’azione circolare. Il feedback loop è il passaggio essenziale: i dati d’uso migliorano le decisioni successive. Se questi moduli fossero app separate avremmo più funzioni ma meno valore. Nella demo seguirò esattamente questo percorso, usando un solo capo come filo narrativo.""",
    """La UI è stata ricostruita per far emergere decisioni, non soltanto numeri. La dashboard mette in primo piano la CO₂ equivalente stimata, ma nello stesso contesto mostra stato, gap e rotazione. “Stato” risponde a cosa richiede cura; “gap” a cosa manca davvero; “rotazione” a quali capi meritano di rientrare in un outfit. La dicitura in basso è deliberata: questa è una mappa delle superfici effettive, non uno screenshot inventato. Durante la demo apro le route reali /dashboard, /items/:id e /today.""",
    """Risposta diretta alla domanda sui ghost garments: sì, sono rilevati. La regola operativa è esatta e non usa machine learning. Un capo è ghost se non ha alcun wear registrato, è posseduto da almeno la soglia scelta—30 giorni per default—ed è ancora attivo. La data di acquisto è preferita; in sua assenza si usa la data di inserimento. La UI mostra il conteggio e la lista nella dashboard, il badge nel dettaglio e il bonus di riscoperta negli outfit. Dico “mai indossato dopo X giorni”, non “non usato di recente”: un capo indossato una volta anni fa non è ghost secondo questa definizione.""",
    """Qui mostro il primo modello addestrato da noi. Non addestriamo Fashion-CLIP da zero: lo congeliamo e lo usiamo come estrattore di 512 feature. Su queste feature alleniamo una testa MLP 512-256-128-3 con circa 165 mila parametri. È transfer learning. Sul test locale di 83 immagini l’accuracy è circa 94 per cento e la F1 della classe danneggiato è 0,981. Ma il limite è importante: lo split è un holdout locale casuale, non una validazione esterna per sorgente. Quindi presento il risultato come prova tecnica del prototipo, non come garanzia su foto reali.""",
    """Il secondo modello lavora sull’intero guardaroba, non sulle foto. Riceve 14 feature aggregate—conteggi per ruolo, stagionalità, colori, neutri, ghost ratio—e produce sei probabilità indipendenti. È multi-label perché uno stesso armadio può avere contemporaneamente pochi capispalla e poca varietà di colori. Le metriche sono alte sul dataset simulato: micro-F1 0,943 e macro-F1 0,932. La precisazione da dire subito è che le etichette derivano da regole esperte con rumore. Oggi la rete dimostra la pipeline e apprende quelle regole; per validare il concetto servono feedback e guardaroba reali.""",
    """Il recommender non è una rete addestrata. Il punteggio combina 55 per cento compatibilità cromatica, 35 per cento meteo, un bonus fino al 15 per cento per capi ghost eleggibili e un piccolo termine di feedback, più o meno 4 per cento. Il feedback quindi personalizza l’ordine, ma non domina il risultato. È importante non ripetere una vecchia descrizione: nel ranking attuale non entrano Chroma né gli embedding Fashion-CLIP. Il vantaggio delle regole qui è la spiegabilità: posso dire perché un outfit è stato proposto e mantenere il comportamento stabile anche senza servizi esterni.""",
    """La dashboard verde non misura emissioni fisiche. Calcola una stima di scenario: media indicativa della categoria moltiplicata per un fattore associato all’azione. L’esempio interno mostra jeans da 32 chilogrammi CO₂ equivalente per 70 per cento in caso di riparazione, cioè 22,4. Questi numeri servono a dare ordine di grandezza e motivazione, ma non sono una LCA per materiale e dipendono dall’ipotesi che riuso o riparazione sostituiscano davvero un nuovo acquisto. All’esame uso sempre le parole “stima” e “CO₂ equivalente”, e propongo come sviluppo range d’incertezza e LCA per materiale.""",
    """Questa slide riassume il confine tecnico. Fashion-CLIP è pre-addestrato da altri e viene usato in inference. Le due MLP—stato e gap—sono addestrate da noi e hanno checkpoint riproducibili. Ghost, outfit, cost-per-wear e stima CO₂ sono regole o tabelle. LLM e diffusion sono modelli generativi esterni e opzionali. Questa tassonomia evita l’errore di chiamare AI qualsiasi calcolo. Il criterio che abbiamo seguito è semplice: addestriamo un modello soltanto quando c’è un’incertezza che una formula esplicita non risolve meglio.""",
    """L’architettura spiega perché il progetto può evolvere. Il client React è separato dalle API FastAPI sotto /api/v1. I servizi contengono la logica di wear, statistiche, gap, recommender e circolarità. Dati e checkpoint sono separati ancora: SQLite e filesystem per l’MVP, PyTorch per le due teste e Fashion-CLIP per la percezione. LLM e try-on si attivano solo se configurati. Oggi il sistema è single-user e locale: è una scelta di prototipo e privacy. Per la produzione servono autenticazione, isolamento per utente e storage remoto sicuro.""",
    """Per l’esame ho separato ogni fitting in un notebook autonomo ed eseguito. I primi due ricostruiscono i modelli caricati dal prodotto. Gli altri tre—rischio ghost con logistic regression, previsione wear con random forest e cluster con K-means—sono esperimenti accademici e non sono collegati alla UI. Questa distinzione è importante: il vecchio notebook unico aveva un esperimento ghost con pochissimi positivi e non era difendibile. Il nuovo notebook genera un dataset didattico non degenere, riporta ROC-AUC 0,713 e dichiara apertamente che il risultato è sintetico.""",
    """Sì, possiamo creare un’app per smartphone senza riscrivere il backend. Oggi l’interfaccia è già responsive e l’upload usa la camera posteriore del browser. La via più breve è una PWA: manifest, service worker e coda offline. Se serve una pubblicazione rapida sugli store mantenendo la UI web, Capacitor è il passaggio intermedio. Expo/React Native ha più senso soltanto se camera nativa, push affidabili, offline-first e UX mobile diventano il centro del prodotto; riutilizzeremmo API e tipi, ma ricostruiremmo i componenti. Prima di tutto servono auth, HTTPS e governance delle immagini.""",
    """Questa non è una slide difensiva: è il piano di validazione. Per generalizzare la diagnosi dello stato dobbiamo fare split per sorgente e un test esterno. Per i gap dobbiamo sostituire o integrare le label simulate con feedback reale e audit dei bias. Per l’impatto servono LCA per materiale e intervalli d’incertezza. Per la privacy servono autenticazione, cifratura, retention e una valutazione GDPR/DPIA perché le foto possono includere corpo e ambiente domestico. Infine SQLite e filesystem vanno sostituiti con Postgres, object storage e job asincroni quando passiamo al multiutente.""",
    """Questa è la scaletta della demo, da provare prima dell’esame con le tab già aperte. In tre minuti e mezzo mostro una singola storia: guardaroba e dettaglio, wear log, ghost con soglia, diagnosi dello stato, gap multi-label e azione circolare che aggiorna l’impatto. Non apro tutte le pagine e non provo ogni pulsante. Per ogni schermata uso sempre la stessa formula orale: qual è l’input, quale trasformazione avviene, qual è l’output, perché aiuta l’utente e qual è il limite. Se una chiamata generativa non è configurata, lo dico e mostro il fallback.""",
    """Chiudo con tre messaggi. Primo: il prodotto trasforma uso, stato, gap e ghost in decisioni concrete che possono allungare la vita dei capi. Secondo: abbiamo addestrato due reti e dichiarato con precisione tutto ciò che è pre-addestrato, generativo o rule-based. Terzo: la qualità del progetto non dipende soltanto dalla metrica migliore, ma dalla riproducibilità dei notebook, dalla visibilità dei limiti e da una roadmap realistica per dati, privacy e mobile. La frase finale è: compra meno, usa meglio, dichiara l’incertezza. Poi mi fermo e apro alle domande.""",
]


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()

    deck = Presentation(args.input)
    if len(deck.slides) != len(NOTES):
        raise ValueError(f"Expected {len(NOTES)} slides, found {len(deck.slides)}")

    for slide, note in zip(deck.slides, NOTES, strict=True):
        frame = slide.notes_slide.notes_text_frame
        frame.clear()
        frame.text = note

    deck.core_properties.title = "ClosetAI — Esame universitario ufficiale"
    deck.core_properties.subject = "Prodotto AI per uso, cura e circolarità del guardaroba"
    deck.core_properties.comments = "Include note del relatore in italiano su tutte le slide."
    args.output.parent.mkdir(parents=True, exist_ok=True)
    deck.save(args.output)
    print(f"Saved {args.output} with {len(NOTES)} speaker notes")


if __name__ == "__main__":
    main()
